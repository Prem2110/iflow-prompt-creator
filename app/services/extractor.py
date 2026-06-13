import base64
import csv
import io
import logging
from pathlib import Path

logger = logging.getLogger(__name__)

_MIME_MAP = {
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".gif": "image/gif",
    ".webp": "image/webp",
}


async def extract(file, data: bytes) -> list[dict]:
    """
    Returns a list of extracted items. Most file types yield one item;
    diagram-heavy PDFs yield one image item per page.

    Each item is one of:
      {"kind": "text",  "name": str, "content": str}
      {"kind": "image", "name": str, "mime": str, "b64": str}
    """
    suffix = Path(file.filename or "").suffix.lower()
    logger.info("Extracting %s (%d bytes)", file.filename, len(data))

    if suffix == ".pdf":
        return _pdf(file.filename, data)
    if suffix in (".docx", ".doc"):
        return [{"kind": "text", "name": file.filename, "content": _docx(data)}]
    if suffix == ".pptx":
        return [{"kind": "text", "name": file.filename, "content": _pptx(data)}]
    if suffix in (".xlsx", ".xls"):
        return [{"kind": "text", "name": file.filename, "content": _xlsx(data)}]
    if suffix == ".csv":
        return [{"kind": "text", "name": file.filename, "content": _csv(data)}]
    if suffix in (".txt", ".json", ".yaml", ".yml", ".xml", ".wsdl"):
        return [{"kind": "text", "name": file.filename, "content": data.decode("utf-8", errors="replace")}]
    if suffix in _MIME_MAP:
        return [{
            "kind": "image",
            "name": file.filename,
            "mime": _MIME_MAP[suffix],
            "b64": base64.b64encode(data).decode(),
        }]

    # Fallback: try utf-8 text
    logger.warning("Unknown extension '%s' for '%s', falling back to plain text", suffix, file.filename)
    return [{"kind": "text", "name": file.filename, "content": data.decode("utf-8", errors="replace")}]


# PDFs with fewer than this many chars per page on average are treated as
# diagram-heavy and rendered as images instead of extracted as text.
_PDF_DIAGRAM_THRESHOLD = 1200
# Never render more than this many pages as images (token cost control).
_PDF_MAX_IMAGE_PAGES = 10
# Render resolution in DPI.
_PDF_RENDER_DPI = 150


def _pdf(name: str, data: bytes) -> list[dict]:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    pages_text = [page.extract_text() or "" for page in reader.pages]
    total_pages = len(pages_text)
    avg_chars = sum(len(t) for t in pages_text) / max(total_pages, 1)

    if avg_chars < _PDF_DIAGRAM_THRESHOLD:
        logger.info(
            "%s: avg %.0f chars/page < %d — rendering as images",
            name, avg_chars, _PDF_DIAGRAM_THRESHOLD,
        )
        return _pdf_as_images(name, data, total_pages)

    logger.info("%s: avg %.0f chars/page — using text extraction", name, avg_chars)
    return [{"kind": "text", "name": name, "content": "\n".join(pages_text)}]


def _pdf_as_images(name: str, data: bytes, total_pages: int) -> list[dict]:
    import fitz  # pymupdf

    doc = fitz.open(stream=data, filetype="pdf")
    items: list[dict] = []
    pages_to_render = min(total_pages, _PDF_MAX_IMAGE_PAGES)

    for i in range(pages_to_render):
        pix = doc[i].get_pixmap(dpi=_PDF_RENDER_DPI)
        img_bytes = pix.tobytes("png")
        label = f"{name} — page {i + 1}" if total_pages > 1 else name
        items.append({
            "kind": "image",
            "name": label,
            "mime": "image/png",
            "b64": base64.b64encode(img_bytes).decode(),
        })

    doc.close()

    if total_pages > _PDF_MAX_IMAGE_PAGES:
        items.append({
            "kind": "text",
            "name": f"{name} (truncated)",
            "content": (
                f"[Note: Only the first {_PDF_MAX_IMAGE_PAGES} of {total_pages} pages "
                f"were sent as images. The remaining pages were omitted.]"
            ),
        })

    return items


def _docx(data: bytes) -> str:
    try:
        from docx import Document
        doc = Document(io.BytesIO(data))
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception:
        # Old Word 97-2003 binary .doc — fall back to OLE2 text extraction
        return _doc_legacy(data)


def _doc_legacy(data: bytes) -> str:
    """Extract text from old .doc files (MHTML/Confluence export, OLE2 binary, RTF)."""
    import re

    # MHTML / Confluence "Export to Word" — MIME multipart with embedded HTML
    if data[:4] in (b"Date", b"MIME", b"Cont") or b"multipart/related" in data[:512]:
        return _mhtml(data)

    # RTF saved as .doc (magic: {\rtf)
    if data[:5] == b"{\\rtf":
        rtf = data.decode("latin-1", errors="replace")
        text = re.sub(r"\\[a-z]+\-?\d* ?", "", rtf)
        text = re.sub(r"[{}\\]", "", text)
        return re.sub(r"\s{3,}", "\n", text).strip()

    # OLE2 compound document (magic: D0 CF 11 E0)
    if data[:4] == b"\xd0\xcf\x11\xe0":
        try:
            import olefile
            ole = olefile.OleFileIO(io.BytesIO(data))
            all_text = []
            for stream_name in ("WordDocument", "1Table", "0Table"):
                if not ole.exists(stream_name):
                    continue
                raw = ole.openstream(stream_name).read()
                for m in re.finditer(rb"(?:[\x20-\x7e][\x00]){5,}", raw):
                    chunk = m.group(0).decode("utf-16-le", errors="ignore").strip()
                    if len(chunk) > 4:
                        all_text.append(chunk)
            ole.close()
            if all_text:
                return "\n".join(all_text)
        except Exception as exc:
            raise ValueError(f"Failed to extract text from .doc file: {exc}") from exc

    raise ValueError(
        "Cannot read this .doc file. Please convert it to .docx (File → Save As → Word Document) and re-upload."
    )


def _mhtml(data: bytes) -> str:
    """Extract text from an MHTML file (multipart MIME with embedded HTML)."""
    import email
    import quopri
    from lxml import html as lhtml

    msg = email.message_from_bytes(data)
    html_content = None

    for part in msg.walk():
        ct = part.get_content_type()
        if ct == "text/html":
            payload = part.get_payload(decode=False)
            cte = part.get("Content-Transfer-Encoding", "").lower()
            if isinstance(payload, str):
                if cte == "quoted-printable":
                    payload = quopri.decodestring(payload.encode("latin-1")).decode("utf-8", errors="replace")
                html_content = payload
                break

    if not html_content:
        raise ValueError("No HTML part found in MHTML file.")

    tree = lhtml.fromstring(html_content)
    # Remove script/style noise
    for tag in tree.xpath("//script | //style"):
        tag.getparent().remove(tag)

    block_tags = {"p", "li", "td", "th", "h1", "h2", "h3", "h4", "h5", "h6", "tr", "div"}
    lines = []
    for el in tree.iter():
        if not isinstance(el.tag, str):
            continue  # skip comments, PIs
        if el.tag in block_tags:
            block = el.text_content().strip()
            if block:
                lines.append(block)

    seen = set()
    unique = []
    for line in lines:
        if line not in seen:
            seen.add(line)
            unique.append(line)
    return "\n".join(unique)


def _pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
    return "\n".join(lines)


def _xlsx(data: bytes) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    lines = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines.append(f"=== Sheet: {sheet_name} ===")
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join(str(c) if c is not None else "" for c in row)
            if row_text.strip():
                lines.append(row_text)
    return "\n".join(lines)


def _csv(data: bytes) -> str:
    text = data.decode("utf-8", errors="replace")
    lines = []
    reader = csv.reader(io.StringIO(text))
    for row in reader:
        lines.append(", ".join(row))
    return "\n".join(lines)
